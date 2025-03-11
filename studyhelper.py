import os
import streamlit as st
from io import BytesIO
from dotenv import load_dotenv

# OpenAI 모듈이 없어도 Gemini API로 대체할 수 있도록 설정
try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

from pathlib import Path
import docx2txt
import pdfplumber
from pptx import Presentation
import random
import subprocess
import nltk
from nltk.corpus import stopwords
import google.generativeai as genai
import PIL.Image
import requests
from concurrent.futures import ThreadPoolExecutor
import datetime
from PIL import Image
import io
from google.generativeai import types  # 올바른 모듈 경로 사용
from docx import Document
import tempfile  # 임시 파일 처리를 위한 모듈

###############################################################################
# 환경 변수 로드 및 설정
###############################################################################
dotenv_path = ".env"
load_dotenv(dotenv_path=dotenv_path)
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
USE_GEMINI_ALWAYS = os.getenv("USE_GEMINI_ALWAYS", "False").lower() == "true"

# Gemini API 설정
genai.configure(api_key=GEMINI_API_KEY)

###############################################################################
# Streamlit 및 NLTK 설정
###############################################################################
nltk_data_dir = "/tmp/nltk_data"
os.makedirs(nltk_data_dir, exist_ok=True)
nltk.data.path.append(nltk_data_dir)

@st.cache_data(show_spinner=False)
def get_stopwords():
    try:
        sw = stopwords.words("english")
    except LookupError:
        nltk.download("stopwords", download_dir=nltk_data_dir)
        sw = stopwords.words("english")
    return set(sw)

english_stopwords = get_stopwords()
korean_stopwords = {"이", "그", "저", "것", "수", "등", "들", "및", "더"}
final_stopwords = english_stopwords.union(korean_stopwords)

###############################################################################
# 병렬 제한 (성능 개선)
###############################################################################
MAX_WORKERS = 2

###############################################################################
# OpenAI & Gemini 설정
###############################################################################
if not OPENAI_API_KEY or OpenAI is None or USE_GEMINI_ALWAYS:
    use_gemini_always = True
    openai_client = None
else:
    use_gemini_always = False
    # base_url은 예시. 필요 시 "https://generativelanguage.googleapis.com/v1beta/openai/"
    openai_client = OpenAI(api_key=OPENAI_API_KEY, base_url="https://api.openai.com/v1")

###############################################################################
# Gemini (캐싱 예시)
###############################################################################
@st.cache_data(show_spinner=False)
def ask_gemini_cached(prompt, temperature=0.2):
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        # prompt를 키워드 인자가 아니라 위치 인자로 전달
        response = model.generate_content(
            prompt,
            generation_config={"temperature": temperature}
        )
        return response.text.strip()
    except Exception as e:
        return f"Gemini API 호출 오류: {e}"

###############################################################################
# OpenAI 호출 실패 시 Gemini로 fallback
###############################################################################
@st.cache_data(show_spinner=False)
def ask_gpt(messages, model_name="gpt-4", temperature=0.2):
    """OpenAI GPT 호출. 실패 시 Gemini API로 fallback."""
    if use_gemini_always or not openai_client:
        return _ask_gemini(messages, temperature=temperature)
    else:
        try:
            resp = openai_client.chat.completions.create(
                model=model_name,
                messages=messages,
                temperature=temperature,
            )
            return resp.choices[0].message.content.strip()
        except Exception as e:
            st.warning(f"OpenAI 호출 실패: {e}, Gemini로 전환합니다.")
            return _ask_gemini(messages, temperature=temperature)

def _ask_gemini(messages, temperature=0.2):
    try:
        system_msg = next((m["content"] for m in messages if m["role"]=="system"), "")
        user_msg = next((m["content"] for m in messages if m["role"]=="user"), "")
        prompt = f"{system_msg}\n\n사용자 질문: {user_msg}"
        return ask_gemini_cached(prompt, temperature=temperature)
    except Exception as e:
        st.error(f"Gemini API fallback 오류: {e}")
        return "AI 호출 중 오류 발생"

###############################################################################
# 이미지 리사이즈 후 분석 (간단 예시)
###############################################################################
@st.cache_data(show_spinner=False)
def analyze_image_resized(file_bytes, max_size=(800, 800)):
    try:
        image = PIL.Image.open(io.BytesIO(file_bytes))
        image.thumbnail(max_size)
        buffer = io.BytesIO()
        image.save(buffer, format='JPEG')
        resized_image_bytes = buffer.getvalue()

        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(
            [types.Part.from_bytes(resized_image_bytes, mime_type='image/jpeg')],
            generation_config={"temperature": 0.2}
        )
        return response.text.strip()
    except Exception as e:
        return f"이미지 분석 오류: {e}"

###############################################################################
# 문서 파싱 (PDF, DOCX, PPT)
###############################################################################
@st.cache_data(show_spinner=False)
def parse_docx(file_bytes):
    try:
        # BytesIO 객체로 전달된 내용을 임시 DOCX 파일로 저장
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            content = file_bytes.getvalue() if hasattr(file_bytes, "getvalue") else file_bytes
            tmp.write(content)
            tmp_path = tmp.name
        text = docx2txt.process(tmp_path)
        os.remove(tmp_path)  # 임시 파일 삭제
        return text
    except Exception as e:
        return f"📄 DOCX 파일 분석 오류: {e}"

@st.cache_data(show_spinner=False)
def parse_pdf(file_bytes):
    text_list = []
    try:
        with pdfplumber.open(file_bytes) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                text_list.append(text)
        return "\n".join(text_list)
    except Exception:
        return "📄 PDF 파일 분석 오류"

@st.cache_data(show_spinner=False)
def parse_ppt(file_bytes):
    text_list = []
    try:
        prs = Presentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_list.append(shape.text)
        return "\n".join(text_list)
    except Exception:
        return "📄 PPTX 파일 분석 오류"

def analyze_file(file_info):
    ext = file_info["ext"]
    data = file_info["data"]
    if ext == "docx":
        return parse_docx(BytesIO(data))
    elif ext == "pdf":
        return parse_pdf(BytesIO(data))
    elif ext == "pptx":
        return parse_ppt(BytesIO(data))
    elif ext in ["png", "jpg", "jpeg"]:
        return analyze_image_resized(data)
    else:
        return "❌ 지원하지 않는 파일 형식입니다."

###############################################################################
# 여러 파일 병합 (병렬 처리)
###############################################################################
@st.cache_data(show_spinner=False)
def merge_documents(file_list):
    def process_file(f):
        bytes_data = f.getvalue()
        return analyze_file({
            "ext": f.name.split(".")[-1].lower(),
            "data": bytes_data
        })

    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        results = list(executor.map(process_file, file_list))

    return "\n\n".join(results)

###############################################################################
# Gemini 캐시 관리 함수
###############################################################################
def manage_gemini_cache():
    from google.generativeai import types
    try:
        cache_client = genai.Client(api_key=GEMINI_API_KEY)
        model_name = "models/gemini-1.5-flash-002"

        # 캐시 생성
        cache = cache_client.caches.create(
            model=model_name,
            config=types.CreateCachedContentConfig(contents=['hello'])
        )
        st.write("캐시 생성됨:", cache)

        # 캐시 TTL 업데이트 (2시간)
        ttl_seconds = int(datetime.timedelta(hours=2).total_seconds())
        cache_client.caches.update(
            name=cache.name,
            config=types.UpdateCachedContentConfig(ttl=f"{ttl_seconds}s")
        )
        st.write("캐시 TTL이 2시간으로 업데이트됨.")

        # 캐시 삭제
        cache_client.caches.delete(name=cache.name)
        st.write("캐시 삭제됨.")
    except Exception as e:
        st.error(f"Gemini 캐시 관리 에러: {e}")

###############################################################################
# 커뮤니티 문서 구조 예시 (버전관리)
# 예시:
# st.session_state.community_posts = [
#     {
#       "id": 1,
#       "title": "법률문서 초안",
#       "content": "여기에 문서 본문이 들어감...",
#       "owner": "익명_101",
#       "history": [
#           {
#             "user": "익명_235",
#             "time": "2025-03-11 14:23",
#             "content": "수정된 문서 내용..."
#           }
#       ]
#     },
#     ...
# ]
###############################################################################

def main():
    st.title("📚 Thinkhelper")
    st.markdown("""
**Thinkhelper**는 AI 기반으로 파일(문서/이미지)을 빠르게 분석하고, 자유롭게 대화할 수 있도록 설계되었습니다.

- 파일을 업로드하면 분석 결과(요약, 수정 제안 등)를 볼 수 있습니다.
- 파일 없이도 자유로운 AI 대화가 가능합니다.
- 추가로, 새 문서를 AI로 자동 생성할 수 있는 기능을 제공합니다.
- **커뮤니티 탭에서는 익명 게시글을 올리고, 다른 사용자가 문서를 수정**(버전관리)할 수 있습니다.
""")

    if st.sidebar.button("Gemini Cache 관리"):
        manage_gemini_cache()

    tab = st.sidebar.radio("🔎 메뉴 선택", ("파일 분석 & GPT 채팅", "AI 문서 생성", "커뮤니티"))

    if "community_posts" not in st.session_state:
        # 초기 예시
        st.session_state.community_posts = []

    if tab == "파일 분석 & GPT 채팅":
        st.info("파일(문서, 이미지)을 업로드하면 분석. 업로드 없이도 AI와 대화.")
        uploaded_files = st.file_uploader(
            "📎 파일 업로드 (PDF, PPTX, DOCX, 이미지)",
            type=["pdf", "pptx", "docx", "png", "jpg", "jpeg"],
            accept_multiple_files=True
        )

        if "analysis_result" not in st.session_state:
            st.session_state.analysis_result = ""

        if uploaded_files:
            with st.spinner("파일 분석 중..."):
                analysis = merge_documents(uploaded_files)
                st.session_state.analysis_result = analysis

            st.subheader("📌 분석 결과")
            st.write(st.session_state.analysis_result)

        st.subheader("💬 GPT와 대화하기")
        user_input = st.text_input("질문을 입력하세요")
        if st.button("전송"):
            if user_input:
                prompt_context = st.session_state.analysis_result
                messages = [
                    {"role": "system", "content": f"파일 내용: {prompt_context}"},
                    {"role": "user", "content": user_input}
                ]
                response = ask_gpt(messages, model_name="gpt-4", temperature=0.2)
                st.write("AI:", response)

    elif tab == "AI 문서 생성":
        st.header("📝 새 문서 생성 (실시간 AI 지원)")
        st.markdown("문서를 작성하는 동시에 AI가 실시간으로 도움을 드립니다.")
        # 문서 작성 및 API 결과 저장을 위한 초기값 설정
        if "doc_text" not in st.session_state:
            st.session_state.doc_text = ""
        if "api_result" not in st.session_state:
            st.session_state.api_result = ""

        def update_document():
            user_input = st.session_state.doc_text
            messages = [
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": user_input}
            ]
            # 사용자가 작성한 문서를 기반으로 API 호출
            result = ask_gpt(messages, model_name="gpt-4", temperature=0.2)
            st.session_state.api_result = result

        # on_change 콜백을 통해 사용자가 내용을 작성할 때마다 update_document 함수가 호출됨
        st.text_area("✏️ 문서 작성", key="doc_text", height=400, on_change=update_document)
        st.subheader("AI 도움 결과")
        st.write(st.session_state.get("api_result", ""))

        doc_type = st.selectbox("문서 형식 선택", ["DOCX (워드 문서)", "텍스트 파일 (.txt)"])
        if st.button("📥 문서 다운로드"):
            if doc_type.startswith("DOCX"):
                doc = Document()
                doc.add_paragraph(st.session_state.doc_text)
                buffer = BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                st.download_button(
                    label="DOCX 다운로드",
                    data=buffer,
                    file_name="generated_document.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
            else:
                st.download_button(
                    label="TXT 다운로드",
                    data=st.session_state.doc_text,
                    file_name="generated_document.txt",
                    mime="text/plain"
                )

    else:
        # 커뮤니티 탭 (버전 관리: 누가 수정했는지 기록)
        st.info("익명으로 게시글(문서) 등록 후, 다른 사용자가 문서를 수정하면 버전 관리 이력을 남깁니다.")

        search_query = st.text_input("🔍 검색 (제목 또는 내용)")
        st.subheader("새 게시글(문서) 작성")
        title = st.text_input("제목")
        content = st.text_area("본문")
        files_uploaded = st.file_uploader("📎 파일 업로드 (선택)", type=["pdf", "pptx", "docx", "png", "jpg", "jpeg"], accept_multiple_files=True)

        if st.button("게시글 등록"):
            if title.strip() and content.strip():
                post_files = []
                if files_uploaded:
                    for uf in files_uploaded:
                        ext = uf.name.split(".")[-1].lower()
                        file_bytes = uf.getvalue()
                        post_files.append({"name": uf.name, "ext": ext, "data": file_bytes})
                new_id = len(st.session_state.community_posts) + 1
                new_post = {
                    "id": new_id,
                    "title": title,
                    "content": content,
                    "owner": f"익명_{random.randint(100,999)}",
                    "files": post_files,
                    "history": []  # 수정 이력
                }
                st.session_state.community_posts.append(new_post)
                st.success("게시글(문서)이 등록되었습니다.")
            else:
                st.error("제목과 내용을 모두 입력해야 합니다.")

        st.subheader("📜 게시글(문서) 목록")
        for idx, post in enumerate(st.session_state.community_posts):
            # 검색
            if (not search_query) or (search_query.lower() in post["title"].lower() or search_query.lower() in post["content"].lower()):
                with st.expander(f"{idx+1}. {post['title']}"):
                    st.write(post["content"])

                    # 파일 목록 표시
                    if post.get("files"):
                        st.markdown("**첨부파일 목록**")
                        for f_idx, f_info in enumerate(post["files"]):
                            st.write(f"- {f_info['name']}")

                    # 수정하기 버튼
                    edit_key = f"edit_{post['id']}"
                    if st.button("수정하기", key=f"edit_btn_{post['id']}"):
                        st.session_state[edit_key] = True

                    if edit_key in st.session_state and st.session_state[edit_key]:
                        st.mar
